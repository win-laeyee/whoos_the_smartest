import time
from typing import Dict, Any
import logging
import os

from moviepy.editor import VideoFileClip
import fitz
import docx
from pptx import Presentation

from backend.src.utils.app_init import init_gemini_llm
from backend.src.utils.helper import check_file_type
from backend.src.utils.constants import VIDEO, WORD_DOCUMENT, PDF_DOCUMENT, IMAGE, PPT_SLIDE

import google.generativeai as genai
from google.generativeai.types import File


def generate_notes(file_path: str):
    model = init_gemini_llm()

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found at {file_path}")
    
    file_type = check_file_type(file_path)

    if file_type in [VIDEO, IMAGE]:
        file = upload_file(file_path, file_type)
        notes = get_notes_from_media(file, model)
        cleanup_file(file)
    elif file_type in [PDF_DOCUMENT, WORD_DOCUMENT, PPT_SLIDE]:
        extracted_text = extract_text(file_path, file_type)
        notes = get_notes_from_document(extracted_text, model)
    
    return notes
   

def cleanup_file(file):
    genai.delete_file(file.name)
    logging.info(f'Deleted file {file.uri}')


def upload_file(file_path: str, file_type: str):
    if file_type == VIDEO:
        return upload_video_file(file_path)
    elif file_type == IMAGE:
        return upload_image_file(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")


def upload_video_file(video_path: str) -> File:

    video_metadata = get_video_metadata(video_path)
    logging.info(f"metadata: {video_metadata}")

    if video_metadata['duration'] >= 7200:
        logging.info(f"Duration of the video is {video_metadata['duration']}")
        raise ValueError("Video file is too long. Make sure it does not exceed 2 hours.")
        
    video_file_name = video_path
    logging.info(f"Uploading file...")
    video_file = genai.upload_file(path=video_file_name)
    logging.info(f"Completed upload: {video_file}")
    print(f"Completed upload: {video_file}")

    while video_file.state.name == "PROCESSING":
        print('.', end='')
        time.sleep(10)
        video_file = genai.get_file(video_file.name)

    if video_file.state.name == "FAILED":
        raise ValueError(video_file.state.name)
    
    return video_file


def get_video_metadata(video_path: str) -> Dict[Any, Any]:
    try:
        clip = VideoFileClip(video_path)
        video_metadata = {
            'duration': clip.duration,
            'width': clip.w,
            'height': clip.h,
            'fps': clip.fps
        }
        return video_metadata
    except Exception as e:
        print(f"An error occurred: {e}")
        return {}
    

def upload_image_file(image_path):
    image_file = genai.upload_file(path=image_path)
    logging.info(f"Completed upload: {image_file}")
    print(f"Completed upload: {image_file}")

    return image_file






def get_notes_from_media(media_file: File, model) -> str:
    prompt = "Generate a summary of notes from the media file provided."
    response = model.generate_content([media_file, prompt],
                                request_options={"timeout": 600})
    # response = model.generate_content('Tell me a story about a magic backpack')
    # llm = ChatGoogleGenerativeAI(model = 'gemini-1.5-pro', google_api_key=GOOGLE_API_KEY, temperature = 0)
    # response = llm.invoke("Write me a ballad about LangChain")
    # return response.content

    return response.text



def get_notes_from_document(extracted_text: str, model) -> str:
    prompt = f"Here is the extracted text: \n{extracted_text}\n\nPlease generate a summary or content based on this text."

    response = model.generate_content(prompt)
    print(f"get notes from doc: {response}")
    return response.text





def extract_text(file_path: str, file_type: str):
    if file_type == PDF_DOCUMENT:
        return extract_text_from_pdf(file_path)
    elif file_type == WORD_DOCUMENT:
        return extract_text_from_word(file_path)
    elif file_type == PPT_SLIDE:
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported document type: {file_type}")


def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text



def extract_text_from_word(docx_path):
    doc = docx.Document(docx_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text



def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text



