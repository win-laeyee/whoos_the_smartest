import docx
import fitz
from pptx import Presentation

from backend.src.utils.constants import PDF_DOCUMENT, PPT_SLIDE, WORD_DOCUMENT


def extract_text_from_pdf(pdf_path: str) -> str:
    """
    Extracts text from a PDF file.

    Args:
        pdf_path (str): The path to the PDF file.

    Returns:
        str: The extracted text from the PDF file.
    """
    doc = fitz.open(pdf_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text


def extract_text_from_word(docx_path: str) -> str:
    """
    Extracts text from a Word document.

    Args:
        docx_path (str): The path to the Word document.

    Returns:
        str: The extracted text from the Word document.
    """
    doc = docx.Document(docx_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text


def extract_text_from_pptx(pptx_path: str) -> str:
    """
    Extracts text from a PowerPoint presentation.

    Args:
        pptx_path (str): The path to the PowerPoint presentation.

    Returns:
        str: The extracted text from the PowerPoint presentation.
    """
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def extract_text(file_path: str, file_type: str) -> str:
    """
    Extracts text from a file based on its type.

    Args:
        file_path (str): The path to the file.
        file_type (str): The type of the file (e.g., 'PDF', 'WORD', 'PPT').

    Returns:
        str: The extracted text from the file.

    Raises:
        ValueError: If the file type is unsupported.
    """
    if file_type == PDF_DOCUMENT:
        return extract_text_from_pdf(file_path)
    elif file_type == WORD_DOCUMENT:
        return extract_text_from_word(file_path)
    elif file_type == PPT_SLIDE:
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported document type: {file_type}")