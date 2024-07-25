from pptx import Presentation
import fitz
import docx
from backend.src.utils.constants import PDF_DOCUMENT, PPT_SLIDE, WORD_DOCUMENT


def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text


def extract_text_from_word(docx_path):
    doc = docx.Document(docx_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text


def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def extract_text(file_path: str, file_type: str):
    if file_type == PDF_DOCUMENT:
        return extract_text_from_pdf(file_path)
    elif file_type == WORD_DOCUMENT:
        return extract_text_from_word(file_path)
    elif file_type == PPT_SLIDE:
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported document type: {file_type}")